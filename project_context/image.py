import os
import base64
import io
from dotenv import load_dotenv
from pydantic.v1 import BaseModel, Field
from typing import List
from pptx import Presentation
import pypdf
from PIL import Image
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.messages import HumanMessage

load_dotenv()


class ImageAnalysis(BaseModel):
    image_index: int = Field(description="Index number of the image being analyzed.")
    description: str = Field(description="Step-by-step description of the diagram or image.")
    type: str = Field(description="Diagram type, e.g., Architecture, User Flow, Data Flow, Sequence, Chart, Mockup.")


class WorkflowReport(BaseModel):
    overall_summary: str = Field(description="High-level workflow across all diagrams.")
    image_analyses: List[ImageAnalysis] = Field(description="Per-image analyses.", min_items=1)


class WorkflowAnalysisAgent:
    def __init__(self):
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY not found in .env file.")

        self.llm = ChatOpenAI(
            model="gpt-4o",
            temperature=0.2,
            openai_api_key=api_key,
        )
        self.parser = JsonOutputParser(pydantic_object=WorkflowReport)
        self.prompt = self._create_prompt()

    def _create_prompt(self):
        prompt_str = """
        You are a system design and process analysis specialist.
        Analyze the images extracted from a presentation. These images can be flowcharts, architecture diagrams, or user journeys.

        1. For each image: describe the process in a step-by-step manner and classify the type of diagram.
        2. After all analyses, create a combined overall summary.
        3. Return ONLY JSON strictly matching this schema:

        {format_instructions}
        """
        return ChatPromptTemplate.from_template(prompt_str)

    def _extract_images_as_base64(self, file_path):
        images = []
        print(f"  -> Extracting images from '{file_path}'...")
        try:
            if file_path.lower().endswith(".pdf"):
                reader = pypdf.PdfReader(file_path)
                for page in reader.pages:
                    if hasattr(page, "images") and page.images:
                        for image_file_object in page.images:
                            img = Image.open(io.BytesIO(image_file_object.data))
                            if img.mode == "RGBA":
                                img = img.convert("RGB")
                            buffered = io.BytesIO()
                            img.save(buffered, format="JPEG")
                            images.append(base64.b64encode(buffered.getvalue()).decode("utf-8"))

            elif file_path.lower().endswith(".pptx"):
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "image"):
                            image_bytes = shape.image.blob
                            img = Image.open(io.BytesIO(image_bytes))
                            if img.mode == "RGBA":
                                img = img.convert("RGB")
                            buffered = io.BytesIO()
                            img.save(buffered, format="JPEG")
                            images.append(base64.b64encode(buffered.getvalue()).decode("utf-8"))
        except Exception as e:
            print(f"  -> Warning: Could not extract images. {e}")

        print(f"  -> Found {len(images)} images.")
        return images

    def analyze_workflows(self, file_path):
        images_base64 = self._extract_images_as_base64(file_path)
        if not images_base64:
            print("  -> No images found to analyze.")
            return None

        prompt_text = self.prompt.format(format_instructions=self.parser.get_format_instructions())

        message_parts = [{"type": "text", "text": prompt_text}]
        for img_data in images_base64:
            message_parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{img_data}"}  # <-- object, not string
            })

        message = HumanMessage(content=message_parts)

        print("  -> Calling OpenAI API for workflow analysis...")
        try:
            response = self.llm.invoke([message])
            raw_content = response.content

            if "```json" in raw_content:
                clean_content = raw_content.split("```json")[1].split("```")[0].strip()
            else:
                clean_content = raw_content

            report_data = self.parser.parse(clean_content)
            print("  -> Analysis complete.")
            return WorkflowReport(**report_data)

        except Exception as e:
            print(f"  -> ERROR during workflow analysis: {e}")
            return None


def display_workflow_report(report: WorkflowReport):
    if not report:
        print("Could not generate a workflow report.")
        return

    print("\n" + "=" * 70)
    print("📊 WORKFLOW AND DIAGRAM ANALYSIS REPORT 📊")
    print("=" * 70)

    print("\n--- 📜 Overall Workflow Summary ---")
    print(report.overall_summary + "\n")

    print("\n--- 🖼️ Detailed Image Analysis ---")
    for analysis in report.image_analyses:
        print(f"\n  ➡️ Image {analysis.image_index} ({analysis.type}):")
        print(f"  {analysis.description}")

    print("\n" + "=" * 70)


if __name__ == "__main__":
    TEAM_FILE_PATH = r"D:\Openai Hackathon Evaluation\ppt\2025CodingPirates.pdf"

    print(f"--- Starting Workflow Analysis for: {TEAM_FILE_PATH} ---")

    if not os.path.exists(TEAM_FILE_PATH):
        print(f"FATAL ERROR: The file '{TEAM_FILE_PATH}' was not found.")
    else:
        agent = WorkflowAnalysisAgent()
        workflow_report = agent.analyze_workflows(TEAM_FILE_PATH)
        if workflow_report:
            display_workflow_report(workflow_report)
